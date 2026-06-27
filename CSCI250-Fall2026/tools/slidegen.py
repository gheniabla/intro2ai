"""
slidegen.py — shared slide generator for CSCI 250 (Fall 2026).

Every weekly deck is built from a simple list of slide dicts so all decks share
one consistent theme. Usage:

    from slidegen import build_deck
    slides = [
        {"type": "title", "title": "Week 1: Python Review", "subtitle": "CSCI 250 — Intro to AI"},
        {"type": "section", "title": "Why Python for AI?"},
        {"type": "bullets", "title": "Topics", "bullets": ["Data types", "Functions", "Modules"]},
        {"type": "code", "title": "Hello, AI", "code": "print('hello')", "caption": "Run in Colab"},
        {"type": "two_col", "title": "Compare", "left_title": "Cloud", "left": ["Claude","Gemini"],
                                                  "right_title": "Local", "right": ["Ollama","HF"]},
        {"type": "closing", "title": "This Week", "bullets": ["Watch videos", "Do Lab A1"]},
    ]
    build_deck(slides, "slides.pptx",
               course="CSCI 250 — Introduction to Artificial Intelligence",
               footer="Palomar College · Fall 2026")

Run `python slidegen.py` to emit a self-test deck (_selftest.pptx).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# ---- Theme -----------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x1F, 0x3A)      # deep navy (Palomar-ish)
TEAL = RGBColor(0x12, 0x8C, 0x8C)      # accent
LIGHT = RGBColor(0xF2, 0xF5, 0xF8)     # light panel
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x6B, 0x72, 0x80)
CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)
CODE_FG = RGBColor(0xE6, 0xE6, 0xE6)

FONT = "Calibri"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)     # 16:9


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h)


def _bar(slide, color, l, t, w, h):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def _set(tf, text, size, color, bold=False, font=FONT, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return p


def _footer(slide, course, footer, idx, total):
    tb = _box(slide, Inches(0.4), Inches(7.05), Inches(9), Inches(0.4))
    _set(tb.text_frame, footer or course, 10, GREY)
    pg = _box(slide, Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.4))
    _set(pg.text_frame, f"{idx}/{total}", 10, GREY, align=PP_ALIGN.RIGHT)


def _title_slide(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, NAVY)
    _bar(slide, TEAL, 0, Inches(3.5), W, Inches(0.08))
    tb = _box(slide, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.6))
    _set(tb.text_frame, s["title"], 40, WHITE, bold=True)
    if s.get("subtitle"):
        sb = _box(slide, Inches(0.9), Inches(3.7), Inches(11.5), Inches(1.2))
        _set(sb.text_frame, s["subtitle"], 22, RGBColor(0xCF, 0xE3, 0xE3))
    return slide


def _section_slide(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, TEAL)
    tb = _box(slide, Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.5))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set(tf, s["title"], 36, WHITE, bold=True)
    return slide


def _content_header(slide, title):
    _bg(slide, WHITE)
    _bar(slide, NAVY, 0, 0, W, Inches(1.15))
    _bar(slide, TEAL, 0, Inches(1.15), W, Inches(0.06))
    tb = _box(slide, Inches(0.6), Inches(0.22), Inches(12.1), Inches(0.8))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set(tf, title, 28, WHITE, bold=True)


def _add_bullets(slide, bullets, left, top, width, height, size=20):
    tb = _box(slide, left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for b in bullets:
        lvl = 0
        text = b
        if isinstance(b, (tuple, list)):
            text, lvl = b[0], b[1]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        r = p.add_run()
        r.text = ("• " if lvl == 0 else "– ") + text
        r.font.size = Pt(size - 2 * lvl)
        r.font.color.rgb = DARK if lvl == 0 else GREY
        r.font.name = FONT
        p.space_after = Pt(8)


def _bullets_slide(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(slide, s["title"])
    _add_bullets(slide, s["bullets"], Inches(0.8), Inches(1.5),
                 Inches(11.7), Inches(5.2), size=s.get("size", 20))
    return slide


def _code_slide(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(slide, s["title"])
    panel = _bar(slide, CODE_BG, Inches(0.6), Inches(1.45),
                 Inches(12.1), Inches(s.get("h", 4.9)))
    tf = panel.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.15)
    lines = s["code"].split("\n")
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = ln if ln else " "
        r.font.size = Pt(s.get("code_size", 15))
        r.font.name = MONO; r.font.color.rgb = CODE_FG
    if s.get("caption"):
        cb = _box(slide, Inches(0.6), Inches(6.45), Inches(11.5), Inches(0.5))
        _set(cb.text_frame, s["caption"], 14, GREY)
    return slide


def _two_col_slide(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(slide, s["title"])
    # left
    lt = _box(slide, Inches(0.8), Inches(1.4), Inches(5.6), Inches(0.6))
    _set(lt.text_frame, s.get("left_title", ""), 20, TEAL, bold=True)
    _add_bullets(slide, s["left"], Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.6), size=18)
    # right
    rt = _box(slide, Inches(6.9), Inches(1.4), Inches(5.6), Inches(0.6))
    _set(rt.text_frame, s.get("right_title", ""), 20, TEAL, bold=True)
    _add_bullets(slide, s["right"], Inches(6.9), Inches(2.0), Inches(5.6), Inches(4.6), size=18)
    return slide


def _closing_slide(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, NAVY)
    tb = _box(slide, Inches(0.9), Inches(0.8), Inches(11.5), Inches(0.9))
    _set(tb.text_frame, s.get("title", "This Week"), 30, WHITE, bold=True)
    _bar(slide, TEAL, Inches(0.9), Inches(1.75), Inches(4), Inches(0.05))
    tb2 = _box(slide, Inches(0.9), Inches(2.1), Inches(11.5), Inches(4.6))
    tf = tb2.text_frame; tf.word_wrap = True
    first = True
    for b in s.get("bullets", []):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = "✓ " + b
        r.font.size = Pt(22); r.font.color.rgb = WHITE; r.font.name = FONT
        p.space_after = Pt(12)
    return slide


_DISPATCH = {
    "title": _title_slide,
    "section": _section_slide,
    "bullets": _bullets_slide,
    "code": _code_slide,
    "two_col": _two_col_slide,
    "closing": _closing_slide,
}


def build_deck(slides, outpath, course="CSCI 250 — Introduction to Artificial Intelligence",
               footer="Palomar College · CSCI 250 · Fall 2026"):
    """Build a .pptx from a list of slide dicts. Returns the output path."""
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    total = len(slides)
    for i, s in enumerate(slides, 1):
        fn = _DISPATCH.get(s["type"])
        if fn is None:
            raise ValueError(f"Unknown slide type: {s['type']}")
        slide = fn(prs, s)
        if s["type"] not in ("title", "section"):
            _footer(slide, course, footer, i, total)
    prs.save(outpath)
    return outpath


if __name__ == "__main__":
    demo = [
        {"type": "title", "title": "Week 0: Slide Engine Self-Test",
         "subtitle": "CSCI 250 — Introduction to Artificial Intelligence · Fall 2026"},
        {"type": "section", "title": "A Section Divider"},
        {"type": "bullets", "title": "A Bullet Slide",
         "bullets": ["Top-level point", ("Sub-point", 1), ("Another sub-point", 1), "Back to top level"]},
        {"type": "code", "title": "A Code Slide",
         "code": "import anthropic\nclient = anthropic.Anthropic()\nprint('hello, AI')",
         "caption": "Runs in Google Colab"},
        {"type": "two_col", "title": "A Two-Column Slide",
         "left_title": "Cloud APIs", "left": ["Claude", "Gemini"],
         "right_title": "Local", "right": ["Ollama", "Hugging Face"]},
        {"type": "closing", "title": "This Week",
         "bullets": ["Review slides", "Run the notebook", "Submit Lab A1"]},
    ]
    out = build_deck(demo, "_selftest.pptx")
    print("wrote", out)
